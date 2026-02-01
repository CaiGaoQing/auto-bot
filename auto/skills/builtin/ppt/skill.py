"""PPT 制作技能"""

from pathlib import Path
from typing import Optional

from auto.core.skill.base import Skill, ToolDefinition
from auto.core.tool.context import ToolContext
from auto.core.tool.result import ToolResult


class PPTSkill(Skill):
    """PPT 制作技能
    
    提供 PPT 创建、编辑、模板应用等功能。
    """
    
    @property
    def name(self) -> str:
        return "ppt"
    
    @property
    def display_name(self) -> str:
        return "PPT 制作"
    
    @property
    def description(self) -> str:
        return "PPT 演示文稿创建、编辑、模板应用、AI 生图"
    
    @property
    def category(self) -> str:
        return "productivity"
    
    @property
    def tools(self) -> list[ToolDefinition]:
        return [
            ToolDefinition(
                name="create_ppt",
                description="创建新的 PPT 演示文稿",
                parameters={
                    "type": "object",
                    "properties": {
                        "title": {
                            "type": "string",
                            "description": "演示文稿标题",
                        },
                        "subtitle": {
                            "type": "string",
                            "description": "副标题",
                        },
                        "output_path": {
                            "type": "string",
                            "description": "输出文件路径",
                        },
                        "template": {
                            "type": "string",
                            "description": "模板名称 (default, business, academic)",
                            "default": "default",
                        },
                    },
                    "required": ["title", "output_path"],
                },
                handler=self.create_ppt,
            ),
            ToolDefinition(
                name="add_slide",
                description="添加幻灯片",
                parameters={
                    "type": "object",
                    "properties": {
                        "ppt_path": {
                            "type": "string",
                            "description": "PPT 文件路径",
                        },
                        "layout": {
                            "type": "string",
                            "enum": ["title", "title_content", "two_column", "blank", "image"],
                            "description": "布局类型",
                            "default": "title_content",
                        },
                        "title": {
                            "type": "string",
                            "description": "幻灯片标题",
                        },
                        "content": {
                            "type": "string",
                            "description": "幻灯片内容 (支持换行符分割要点)",
                        },
                        "notes": {
                            "type": "string",
                            "description": "演讲者备注",
                        },
                    },
                    "required": ["ppt_path", "title"],
                },
                handler=self.add_slide,
            ),
            ToolDefinition(
                name="add_image_slide",
                description="添加带图片的幻灯片",
                parameters={
                    "type": "object",
                    "properties": {
                        "ppt_path": {
                            "type": "string",
                            "description": "PPT 文件路径",
                        },
                        "title": {
                            "type": "string",
                            "description": "幻灯片标题",
                        },
                        "image_path": {
                            "type": "string",
                            "description": "图片文件路径",
                        },
                        "caption": {
                            "type": "string",
                            "description": "图片说明",
                        },
                    },
                    "required": ["ppt_path", "title", "image_path"],
                },
                handler=self.add_image_slide,
            ),
            ToolDefinition(
                name="generate_from_outline",
                description="根据大纲自动生成 PPT",
                parameters={
                    "type": "object",
                    "properties": {
                        "outline": {
                            "type": "string",
                            "description": "PPT 大纲 (Markdown 格式，# 为幻灯片标题)",
                        },
                        "output_path": {
                            "type": "string",
                            "description": "输出文件路径",
                        },
                        "template": {
                            "type": "string",
                            "description": "模板名称",
                            "default": "default",
                        },
                    },
                    "required": ["outline", "output_path"],
                },
                handler=self.generate_from_outline,
            ),
            ToolDefinition(
                name="read_ppt",
                description="读取 PPT 内容",
                parameters={
                    "type": "object",
                    "properties": {
                        "ppt_path": {
                            "type": "string",
                            "description": "PPT 文件路径",
                        },
                    },
                    "required": ["ppt_path"],
                },
                handler=self.read_ppt,
            ),
            ToolDefinition(
                name="export_to_pdf",
                description="将 PPT 导出为 PDF",
                parameters={
                    "type": "object",
                    "properties": {
                        "ppt_path": {
                            "type": "string",
                            "description": "PPT 文件路径",
                        },
                        "output_path": {
                            "type": "string",
                            "description": "PDF 输出路径",
                        },
                    },
                    "required": ["ppt_path"],
                },
                handler=self.export_to_pdf,
            ),
        ]
    
    @property
    def system_prompt(self) -> str:
        return """你是一个专业的 PPT 制作助手，擅长：
- 创建专业的演示文稿
- 设计清晰的内容结构
- 应用合适的布局和模板
- 生成配图和图表

制作原则：
1. 每页内容精简，突出重点
2. 使用项目符号列表
3. 标题简洁有力
4. 保持视觉一致性"""
    
    def _ensure_pptx(self):
        """确保 python-pptx 已安装"""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            return True
        except ImportError:
            raise ImportError("需要安装 python-pptx: pip install python-pptx")
    
    async def create_ppt(
        self,
        ctx: ToolContext,
        title: str,
        output_path: str,
        subtitle: Optional[str] = None,
        template: str = "default",
    ) -> ToolResult:
        """创建新的 PPT"""
        try:
            self._ensure_pptx()
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except ImportError as e:
            return ToolResult.error_result(str(e))
        
        path = Path(output_path).expanduser()
        
        if not ctx.security.is_allowed_path(path):
            return ToolResult.error_result(f"路径不允许: {output_path}")
        
        try:
            # 创建演示文稿
            prs = Presentation()
            
            # 设置幻灯片大小 (16:9)
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            
            # 添加标题幻灯片
            title_slide_layout = prs.slide_layouts[0]  # 标题布局
            slide = prs.slides.add_slide(title_slide_layout)
            
            # 设置标题
            title_shape = slide.shapes.title
            title_shape.text = title
            
            # 设置副标题
            if subtitle and len(slide.placeholders) > 1:
                subtitle_shape = slide.placeholders[1]
                subtitle_shape.text = subtitle
            
            # 保存
            path.parent.mkdir(parents=True, exist_ok=True)
            prs.save(str(path))
            
            return ToolResult.file(
                path=str(path),
                message=f"PPT 已创建: {path.name}",
            )
        except Exception as e:
            return ToolResult.error_result(f"创建 PPT 失败: {str(e)}")
    
    async def add_slide(
        self,
        ctx: ToolContext,
        ppt_path: str,
        title: str,
        layout: str = "title_content",
        content: Optional[str] = None,
        notes: Optional[str] = None,
    ) -> ToolResult:
        """添加幻灯片"""
        try:
            self._ensure_pptx()
            from pptx import Presentation
            from pptx.util import Pt
        except ImportError as e:
            return ToolResult.error_result(str(e))
        
        path = Path(ppt_path).expanduser()
        
        if not ctx.security.is_allowed_path(path):
            return ToolResult.error_result(f"路径不允许: {ppt_path}")
        
        if not path.exists():
            return ToolResult.error_result(f"文件不存在: {ppt_path}")
        
        try:
            prs = Presentation(str(path))
            
            # 选择布局
            layout_map = {
                "title": 0,
                "title_content": 1,
                "two_column": 3,
                "blank": 6,
                "image": 5,
            }
            layout_idx = layout_map.get(layout, 1)
            
            # 确保布局索引有效
            if layout_idx >= len(prs.slide_layouts):
                layout_idx = 1
            
            slide_layout = prs.slide_layouts[layout_idx]
            slide = prs.slides.add_slide(slide_layout)
            
            # 设置标题
            if slide.shapes.title:
                slide.shapes.title.text = title
            
            # 设置内容
            if content and len(slide.placeholders) > 1:
                body_shape = slide.placeholders[1]
                tf = body_shape.text_frame
                tf.clear()
                
                # 分割内容为要点
                points = content.split("\n")
                for i, point in enumerate(points):
                    point = point.strip()
                    if not point:
                        continue
                    
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    
                    # 处理缩进级别
                    if point.startswith("  - ") or point.startswith("    "):
                        p.level = 1
                        point = point.strip("- ").strip()
                    elif point.startswith("- "):
                        p.level = 0
                        point = point[2:]
                    
                    p.text = point
            
            # 添加备注
            if notes:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = notes
            
            # 保存
            prs.save(str(path))
            
            return ToolResult.success_result(
                data={"slide_count": len(prs.slides)},
                message=f"已添加幻灯片，共 {len(prs.slides)} 页",
            )
        except Exception as e:
            return ToolResult.error_result(f"添加幻灯片失败: {str(e)}")
    
    async def add_image_slide(
        self,
        ctx: ToolContext,
        ppt_path: str,
        title: str,
        image_path: str,
        caption: Optional[str] = None,
    ) -> ToolResult:
        """添加带图片的幻灯片"""
        try:
            self._ensure_pptx()
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError as e:
            return ToolResult.error_result(str(e))
        
        ppt_file = Path(ppt_path).expanduser()
        img_file = Path(image_path).expanduser()
        
        if not ctx.security.is_allowed_path(ppt_file):
            return ToolResult.error_result(f"路径不允许: {ppt_path}")
        
        if not ppt_file.exists():
            return ToolResult.error_result(f"PPT 文件不存在: {ppt_path}")
        
        if not img_file.exists():
            return ToolResult.error_result(f"图片文件不存在: {image_path}")
        
        try:
            prs = Presentation(str(ppt_file))
            
            # 使用空白布局
            blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
            slide = prs.slides.add_slide(blank_layout)
            
            # 添加标题
            from pptx.util import Pt
            from pptx.dml.color import RGBColor
            
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(28)
            p.font.bold = True
            
            # 添加图片
            left = Inches(1)
            top = Inches(1.5)
            width = Inches(11)
            
            slide.shapes.add_picture(str(img_file), left, top, width=width)
            
            # 添加说明
            if caption:
                caption_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12), Inches(0.5))
                tf = caption_box.text_frame
                p = tf.paragraphs[0]
                p.text = caption
                p.font.size = Pt(14)
                p.font.italic = True
            
            # 保存
            prs.save(str(ppt_file))
            
            return ToolResult.success_result(
                data={"slide_count": len(prs.slides)},
                message=f"已添加图片幻灯片，共 {len(prs.slides)} 页",
            )
        except Exception as e:
            return ToolResult.error_result(f"添加图片幻灯片失败: {str(e)}")
    
    async def generate_from_outline(
        self,
        ctx: ToolContext,
        outline: str,
        output_path: str,
        template: str = "default",
    ) -> ToolResult:
        """根据大纲生成 PPT"""
        try:
            self._ensure_pptx()
        except ImportError as e:
            return ToolResult.error_result(str(e))
        
        path = Path(output_path).expanduser()
        
        if not ctx.security.is_allowed_path(path):
            return ToolResult.error_result(f"路径不允许: {output_path}")
        
        try:
            # 解析大纲
            slides_data = self._parse_outline(outline)
            
            if not slides_data:
                return ToolResult.error_result("无法解析大纲，请使用 Markdown 格式")
            
            # 创建 PPT
            first_slide = slides_data[0]
            result = await self.create_ppt(
                ctx,
                title=first_slide["title"],
                output_path=output_path,
                subtitle=first_slide.get("subtitle"),
                template=template,
            )
            
            if not result.success:
                return result
            
            # 添加其他幻灯片
            for slide_data in slides_data[1:]:
                await self.add_slide(
                    ctx,
                    ppt_path=output_path,
                    title=slide_data["title"],
                    content=slide_data.get("content", ""),
                    layout="title_content",
                )
            
            return ToolResult.file(
                path=str(path),
                message=f"已根据大纲生成 PPT，共 {len(slides_data)} 页",
            )
        except Exception as e:
            return ToolResult.error_result(f"生成 PPT 失败: {str(e)}")
    
    def _parse_outline(self, outline: str) -> list[dict]:
        """解析 Markdown 大纲"""
        slides = []
        current_slide = None
        content_lines = []
        
        for line in outline.split("\n"):
            line = line.rstrip()
            
            if line.startswith("# "):
                # 保存上一个幻灯片
                if current_slide:
                    current_slide["content"] = "\n".join(content_lines)
                    slides.append(current_slide)
                    content_lines = []
                
                # 新幻灯片
                current_slide = {"title": line[2:].strip()}
            
            elif line.startswith("## "):
                # 副标题或新幻灯片
                if current_slide and not slides:
                    current_slide["subtitle"] = line[3:].strip()
                else:
                    if current_slide:
                        current_slide["content"] = "\n".join(content_lines)
                        slides.append(current_slide)
                        content_lines = []
                    current_slide = {"title": line[3:].strip()}
            
            elif line.startswith("- ") or line.startswith("* "):
                content_lines.append(line[2:].strip())
            
            elif line.strip() and current_slide:
                content_lines.append(line.strip())
        
        # 保存最后一个幻灯片
        if current_slide:
            current_slide["content"] = "\n".join(content_lines)
            slides.append(current_slide)
        
        return slides
    
    async def read_ppt(
        self,
        ctx: ToolContext,
        ppt_path: str,
    ) -> ToolResult:
        """读取 PPT 内容"""
        try:
            self._ensure_pptx()
            from pptx import Presentation
        except ImportError as e:
            return ToolResult.error_result(str(e))
        
        path = Path(ppt_path).expanduser()
        
        if not ctx.security.is_allowed_path(path):
            return ToolResult.error_result(f"路径不允许: {ppt_path}")
        
        if not path.exists():
            return ToolResult.error_result(f"文件不存在: {ppt_path}")
        
        try:
            prs = Presentation(str(path))
            
            slides_data = []
            
            for i, slide in enumerate(prs.slides):
                slide_info = {
                    "index": i + 1,
                    "title": "",
                    "content": [],
                    "notes": "",
                }
                
                # 提取标题
                if slide.shapes.title:
                    slide_info["title"] = slide.shapes.title.text
                
                # 提取内容
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        if shape != slide.shapes.title:
                            slide_info["content"].append(shape.text)
                
                # 提取备注
                if slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text
                    if notes_text:
                        slide_info["notes"] = notes_text
                
                slides_data.append(slide_info)
            
            return ToolResult.success_result(
                data={
                    "file": str(path),
                    "slide_count": len(slides_data),
                    "slides": slides_data,
                },
                message=f"读取 PPT 成功，共 {len(slides_data)} 页",
            )
        except Exception as e:
            return ToolResult.error_result(f"读取 PPT 失败: {str(e)}")
    
    async def export_to_pdf(
        self,
        ctx: ToolContext,
        ppt_path: str,
        output_path: Optional[str] = None,
    ) -> ToolResult:
        """导出为 PDF"""
        # 注意：python-pptx 不直接支持 PDF 导出
        # 需要使用 LibreOffice 或其他工具
        return ToolResult.error_result(
            "PDF 导出需要 LibreOffice，请使用命令: "
            "libreoffice --headless --convert-to pdf your.pptx"
        )
