"""技能包路由 - 完整实现"""

import os
import json
import asyncio
from datetime import datetime
from pathlib import Path
from typing import Optional, List, Any

from fastapi import APIRouter, HTTPException, BackgroundTasks
from pydantic import BaseModel

from auto.core.skill.engine import get_skill_engine
from auto.core.ai.router import get_router

router = APIRouter()

# 工作空间根目录 - 使用项目目录下的 data 文件夹
_PROJECT_ROOT = Path(__file__).parent.parent.parent.parent.parent  # auto/gateway/api/routes -> project root
WORKSPACES_ROOT = _PROJECT_ROOT / "data" / "workspaces"


class ToolExecuteRequest(BaseModel):
    """工具执行请求"""
    tool_name: str
    arguments: dict = {}
    workspace_id: Optional[str] = None
    input_files: List[str] = []  # 输入文件名列表


class SkillExecuteRequest(BaseModel):
    """技能执行请求 - AI 驱动"""
    task: str  # 任务描述，如 "分析这份财报并生成总结"
    workspace_id: str
    input_files: List[str] = []  # 输入文件名列表
    output_format: Optional[str] = None  # 期望的输出格式: pptx, xlsx, pdf, md


class ExecutionResult(BaseModel):
    """执行结果"""
    code: int = 0
    message: str = "success"
    data: dict


# 执行状态存储
execution_status: dict[str, dict] = {}


def get_workspace_paths(workspace_id: str) -> tuple[Path, Path, Path]:
    """获取工作空间路径"""
    ws_path = WORKSPACES_ROOT / workspace_id
    files_path = ws_path / "files"
    outputs_path = ws_path / "outputs"
    return ws_path, files_path, outputs_path


@router.get("/skills")
async def list_skills():
    """列出技能包"""
    engine = get_skill_engine()
    skills = engine.list_skills()
    
    return {
        "code": 0,
        "data": {
            "items": skills,
            "total": len(skills),
        }
    }


@router.get("/skills/{skill_name}")
async def get_skill(skill_name: str):
    """获取技能详情"""
    engine = get_skill_engine()
    skill = engine.get_skill(skill_name)
    
    if not skill:
        return {
            "code": 40401,
            "message": "技能不存在",
        }
    
    return {
        "code": 0,
        "data": {
            "name": skill.name,
            "display_name": skill.display_name,
            "version": skill.version,
            "description": skill.description,
            "category": skill.category,
            "tools": [
                {
                    "name": t.name,
                    "description": t.description,
                    "dangerous": t.dangerous,
                    "parameters": t.parameters if hasattr(t, 'parameters') else {},
                }
                for t in skill.tools
            ],
        }
    }


@router.get("/skills/{skill_name}/tools")
async def get_skill_tools(skill_name: str):
    """获取技能工具列表 (OpenAI 格式)"""
    engine = get_skill_engine()
    tools = engine.get_tools_for_skill(skill_name)
    
    return {
        "code": 0,
        "data": {
            "tools": tools,
        }
    }


@router.post("/skills/{skill_name}/execute")
async def execute_tool(skill_name: str, request: ToolExecuteRequest) -> ExecutionResult:
    """执行技能中的工具"""
    engine = get_skill_engine()
    skill = engine.get_skill(skill_name)
    
    if not skill:
        raise HTTPException(status_code=404, detail="技能不存在")
    
    # 查找工具
    tool = None
    for t in skill.tools:
        if t.name == request.tool_name:
            tool = t
            break
    
    if not tool:
        raise HTTPException(status_code=404, detail="工具不存在")
    
    # 设置工作空间上下文
    context = {}
    if request.workspace_id:
        ws_path, files_path, outputs_path = get_workspace_paths(request.workspace_id)
        if not ws_path.exists():
            raise HTTPException(status_code=404, detail="工作空间不存在")
        
        context["workspace_id"] = request.workspace_id
        context["workspace_path"] = str(ws_path)
        context["files_path"] = str(files_path)
        context["outputs_path"] = str(outputs_path)
        
        # 将输入文件路径添加到参数
        if request.input_files:
            input_file_paths = []
            for filename in request.input_files:
                file_path = files_path / filename
                if file_path.exists():
                    input_file_paths.append(str(file_path))
            context["input_files"] = input_file_paths
    
    try:
        # 执行工具
        result = await engine.execute_tool(
            skill_name=skill_name,
            tool_name=request.tool_name,
            arguments=request.arguments,
            context=context,
        )
        
        # 检查是否有生成的输出文件
        output_files = []
        if request.workspace_id:
            _, _, outputs_path = get_workspace_paths(request.workspace_id)
            for f in outputs_path.glob("*"):
                if f.is_file():
                    output_files.append({
                        "name": f.name,
                        "size": f.stat().st_size,
                        "created_at": datetime.fromtimestamp(f.stat().st_ctime).isoformat(),
                    })
        
        return ExecutionResult(
            data={
                "result": result,
                "output_files": output_files,
            }
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/skills/ai-execute")
async def ai_execute_skill(
    request: SkillExecuteRequest,
    background_tasks: BackgroundTasks,
) -> ExecutionResult:
    """AI 驱动的技能执行 - 根据任务描述自动选择和执行技能"""
    
    # 验证工作空间
    ws_path, files_path, outputs_path = get_workspace_paths(request.workspace_id)
    if not ws_path.exists():
        raise HTTPException(status_code=404, detail="工作空间不存在")
    
    # 生成执行 ID
    execution_id = f"exec_{datetime.now().strftime('%Y%m%d_%H%M%S')}_{os.urandom(4).hex()}"
    
    # 初始化状态
    execution_status[execution_id] = {
        "id": execution_id,
        "status": "pending",
        "task": request.task,
        "workspace_id": request.workspace_id,
        "started_at": datetime.now().isoformat(),
        "progress": 0,
        "message": "准备执行...",
        "output_files": [],
    }
    
    # 在后台执行任务
    background_tasks.add_task(
        _execute_ai_task,
        execution_id,
        request.task,
        request.workspace_id,
        request.input_files,
        request.output_format,
    )
    
    return ExecutionResult(
        data={
            "execution_id": execution_id,
            "status": "pending",
            "message": "任务已提交，正在处理中...",
        }
    )


async def _execute_ai_task(
    execution_id: str,
    task: str,
    workspace_id: str,
    input_files: List[str],
    output_format: Optional[str],
):
    """后台执行 AI 任务"""
    try:
        execution_status[execution_id]["status"] = "running"
        execution_status[execution_id]["progress"] = 10
        execution_status[execution_id]["message"] = "分析任务..."
        
        ws_path, files_path, outputs_path = get_workspace_paths(workspace_id)
        
        # 读取输入文件内容（用于上下文）
        file_contents = {}
        for filename in input_files:
            file_path = files_path / filename
            if file_path.exists():
                # 根据文件类型读取
                suffix = file_path.suffix.lower()
                if suffix in ['.txt', '.md', '.csv', '.json']:
                    try:
                        file_contents[filename] = file_path.read_text(encoding='utf-8')[:10000]  # 限制大小
                    except:
                        file_contents[filename] = f"[二进制文件: {filename}]"
                elif suffix in ['.xlsx', '.xls']:
                    try:
                        import pandas as pd
                        df = pd.read_excel(file_path)
                        file_contents[filename] = df.head(50).to_string()
                    except:
                        file_contents[filename] = f"[Excel 文件: {filename}]"
                elif suffix == '.pdf':
                    try:
                        import pymupdf
                        doc = pymupdf.open(file_path)
                        text = ""
                        for page in doc[:10]:  # 最多读取10页
                            text += page.get_text()
                        file_contents[filename] = text[:10000]
                    except:
                        file_contents[filename] = f"[PDF 文件: {filename}]"
                else:
                    file_contents[filename] = f"[文件: {filename}]"
        
        execution_status[execution_id]["progress"] = 30
        execution_status[execution_id]["message"] = "调用 AI 处理..."
        
        # 构建 AI 提示
        file_context = ""
        if file_contents:
            file_context = "\n\n## 输入文件内容:\n"
            for name, content in file_contents.items():
                file_context += f"\n### {name}:\n```\n{content[:5000]}\n```\n"
        
        output_instruction = ""
        if output_format:
            output_instruction = f"\n\n请将结果以 {output_format} 格式输出。"
        
        prompt = f"""你是一个专业的工作助手。请完成以下任务：

## 任务描述:
{task}
{file_context}
{output_instruction}

请详细分析并给出专业的结果。如果需要生成文件，请明确说明文件内容和格式。
"""
        
        # 调用 AI
        from auto.shared.models import Message, MessageRole
        ai_router = get_router()
        
        response = await ai_router.chat(
            messages=[Message(role=MessageRole.USER, content=prompt)],
        )
        
        ai_result = response.message.content
        
        execution_status[execution_id]["progress"] = 70
        execution_status[execution_id]["message"] = "生成输出文件..."
        
        # 生成输出文件
        output_files = []
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        
        # 根据输出格式生成文件
        if output_format == "xlsx":
            # 生成 Excel
            try:
                output_file = outputs_path / f"result_{timestamp}.xlsx"
                import pandas as pd
                
                # 尝试从 AI 结果中提取表格数据
                # 简单实现：将结果作为单元格内容
                df = pd.DataFrame({"分析结果": [ai_result]})
                df.to_excel(output_file, index=False)
                
                output_files.append({
                    "name": output_file.name,
                    "size": output_file.stat().st_size,
                    "type": "xlsx",
                })
            except Exception as e:
                print(f"生成 Excel 失败: {e}")
        
        elif output_format == "pptx":
            # 生成 PPT
            try:
                from pptx import Presentation
                from pptx.util import Inches, Pt
                
                output_file = outputs_path / f"result_{timestamp}.pptx"
                prs = Presentation()
                
                # 标题页
                slide_layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                title.text = task[:50]
                
                # 内容页（按段落分割）
                paragraphs = ai_result.split('\n\n')
                for i, para in enumerate(paragraphs[:10]):  # 最多10页
                    if para.strip():
                        slide_layout = prs.slide_layouts[1]
                        slide = prs.slides.add_slide(slide_layout)
                        title = slide.shapes.title
                        title.text = f"第 {i+1} 部分"
                        body = slide.shapes.placeholders[1]
                        tf = body.text_frame
                        tf.text = para[:500]
                
                prs.save(output_file)
                
                output_files.append({
                    "name": output_file.name,
                    "size": output_file.stat().st_size,
                    "type": "pptx",
                })
            except Exception as e:
                print(f"生成 PPT 失败: {e}")
        
        elif output_format == "pdf":
            # 生成 PDF
            try:
                from weasyprint import HTML
                
                output_file = outputs_path / f"result_{timestamp}.pdf"
                
                # 将 Markdown 转换为 HTML
                import markdown
                html_content = markdown.markdown(ai_result)
                html = f"""
                <html>
                <head>
                    <meta charset="utf-8">
                    <style>
                        body {{ font-family: Arial, sans-serif; padding: 40px; line-height: 1.6; }}
                        h1, h2, h3 {{ color: #2563EB; }}
                        pre {{ background: #f5f5f5; padding: 10px; border-radius: 5px; }}
                    </style>
                </head>
                <body>
                    <h1>{task[:50]}</h1>
                    {html_content}
                </body>
                </html>
                """
                
                HTML(string=html).write_pdf(output_file)
                
                output_files.append({
                    "name": output_file.name,
                    "size": output_file.stat().st_size,
                    "type": "pdf",
                })
            except Exception as e:
                print(f"生成 PDF 失败: {e}")
        
        # 始终生成 Markdown 结果
        md_file = outputs_path / f"result_{timestamp}.md"
        md_file.write_text(f"# {task}\n\n{ai_result}", encoding="utf-8")
        output_files.append({
            "name": md_file.name,
            "size": md_file.stat().st_size,
            "type": "md",
        })
        
        # 更新状态
        execution_status[execution_id]["status"] = "completed"
        execution_status[execution_id]["progress"] = 100
        execution_status[execution_id]["message"] = "执行完成"
        execution_status[execution_id]["result"] = ai_result[:2000]  # 截断结果
        execution_status[execution_id]["output_files"] = output_files
        execution_status[execution_id]["completed_at"] = datetime.now().isoformat()
        
    except Exception as e:
        execution_status[execution_id]["status"] = "failed"
        execution_status[execution_id]["message"] = str(e)
        execution_status[execution_id]["error"] = str(e)


@router.get("/skills/executions/{execution_id}")
async def get_execution_status(execution_id: str) -> ExecutionResult:
    """获取执行状态"""
    if execution_id not in execution_status:
        raise HTTPException(status_code=404, detail="执行任务不存在")
    
    return ExecutionResult(data=execution_status[execution_id])
