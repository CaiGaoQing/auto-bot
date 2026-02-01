"""对话路由"""

import json
import re
from datetime import datetime
from pathlib import Path
from typing import Optional

from fastapi import APIRouter, HTTPException
from pydantic import BaseModel

from auto.core.ai.router import get_router
from auto.shared.models import Message, MessageRole

router = APIRouter()

# 工作空间根目录
_PROJECT_ROOT = Path(__file__).parent.parent.parent.parent.parent
WORKSPACES_ROOT = _PROJECT_ROOT / "data" / "workspaces"


class ChatRequest(BaseModel):
    """聊天请求"""
    message: str
    model: Optional[str] = None
    workspace_id: Optional[str] = None
    role: Optional[str] = None
    stream: bool = False
    save_to_workspace: bool = True  # 是否保存到工作空间


class ChatResponse(BaseModel):
    """聊天响应"""
    code: int = 0
    message: str = "success"
    data: dict


def _should_save_file(message: str) -> bool:
    """判断是否需要保存文件"""
    message_lower = message.lower()
    
    # 明确需要生成文件的关键词
    save_keywords = [
        # 文档类
        'ppt', '演示', '幻灯片', 'slide', 'powerpoint',
        'excel', '表格', 'xlsx', 'csv',
        '文档', 'document', 'word', 'docx',
        # 代码类
        '写代码', '写脚本', '生成代码', '创建脚本',
        '写一个', '生成一个', '创建一个', '做一个',
        '帮我写', '帮我生成', '帮我创建',
        # 数据类
        'sql', '数据库', 'database', '建表', '表结构',
        'json文件', 'yaml文件', 'xml文件',
        # 输出类
        '导出', 'export', '保存', 'save', '输出文件',
        # PRD/设计
        'prd', '产品文档', '需求文档', '设计文档', '技术方案',
    ]
    
    # 不需要保存文件的关键词（优先级更高）
    no_save_keywords = [
        '什么是', '是什么', '怎么', '如何', '为什么', '解释',
        '分析', '检查', '看看', '理解', '学习',
        '有什么', '哪些', '区别', '对比', '比较',
        '告诉我', '说说', '介绍', '了解',
        '问题', '错误', '修复', 'fix', 'bug',
        '?', '？',
    ]
    
    # 如果包含不保存关键词，返回 False
    if any(kw in message_lower for kw in no_save_keywords):
        # 除非同时有明确的生成关键词
        if not any(kw in message_lower for kw in ['生成', '写', '创建', '做', '导出']):
            return False
    
    # 如果包含保存关键词，返回 True
    if any(kw in message_lower for kw in save_keywords):
        return True
    
    return False


def _detect_file_type(message: str) -> tuple[str, str]:
    """根据消息检测应该生成的文件类型"""
    message_lower = message.lower()
    
    # PPT
    if any(kw in message_lower for kw in ['ppt', '演示', '幻灯片', 'slide', 'powerpoint']):
        return 'pptx', 'ppt'
    
    # Excel
    elif any(kw in message_lower for kw in ['excel', '表格', '工资', '财务', '数据表', 'xlsx', 'csv']):
        return 'xlsx', 'data'
    
    # Java/SpringBoot 项目代码 - 优先于 SQL（因为可能同时包含"数据库"关键词）
    elif any(kw in message_lower for kw in ['springboot', 'spring boot', 'spring-boot', 'java后端', 'java代码', '后端代码', 'controller', 'service', 'entity', '实体类']):
        return 'java_project', 'code'
    
    # SQL - 仅当明确是 SQL 相关时
    elif any(kw in message_lower for kw in ['sql', '建表', '表结构', 'mysql', 'postgresql', 'sqlite']) and not any(kw in message_lower for kw in ['代码', 'code', 'controller', 'service']):
        return 'sql', 'scripts'
    elif '数据库' in message_lower and '设计' in message_lower and not any(kw in message_lower for kw in ['代码', 'code', 'controller']):
        return 'sql', 'scripts'
    
    # Java 单文件
    elif any(kw in message_lower for kw in ['java', '.java']) and 'javascript' not in message_lower:
        return 'java', 'code'
    
    # 其他代码文件
    elif any(kw in message_lower for kw in ['python', '.py']):
        return 'py', 'code'
    elif any(kw in message_lower for kw in ['javascript', '.js']) and 'typescript' not in message_lower:
        return 'js', 'code'
    elif any(kw in message_lower for kw in ['typescript', '.ts', '.tsx']):
        return 'ts', 'code'
    elif any(kw in message_lower for kw in ['go', 'golang', '.go']):
        return 'go', 'code'
    elif any(kw in message_lower for kw in ['rust', '.rs']):
        return 'rs', 'code'
    elif any(kw in message_lower for kw in ['c++', 'cpp', '.cpp']):
        return 'cpp', 'code'
    elif any(kw in message_lower for kw in ['shell', 'bash', '.sh', '脚本']):
        return 'sh', 'scripts'
    
    # 配置文件
    elif any(kw in message_lower for kw in ['json', 'json文件']):
        return 'json', 'data'
    elif any(kw in message_lower for kw in ['yaml', 'yml']):
        return 'yaml', 'data'
    elif any(kw in message_lower for kw in ['xml']):
        return 'xml', 'data'
    
    # HTML
    elif any(kw in message_lower for kw in ['html', '网页', 'webpage']):
        return 'html', 'code'
    
    # 通用代码请求 - 检查是否是代码生成请求
    elif any(kw in message_lower for kw in ['生成代码', '写代码', '编写代码', '后端代码', '前端代码', '代码生成']):
        return 'code_project', 'code'
    
    # 默认 Markdown
    else:
        return 'md', 'docs'


def _extract_code_blocks(content: str) -> list[dict]:
    """从 AI 输出中提取所有代码块"""
    code_blocks = []
    
    # 匹配 ```language:filepath 或 ```language filepath 或 ```language 格式的代码块
    pattern = r'```(\w+)(?::([^\n`]+))?\n(.*?)```'
    matches = re.findall(pattern, content, re.DOTALL)
    
    for lang, filepath, code in matches:
        filepath = filepath.strip() if filepath else None
        code_blocks.append({
            'language': lang,
            'filepath': filepath,
            'code': code.strip()
        })
    
    return code_blocks


def _infer_filepath_from_code(code: str, lang: str) -> Optional[str]:
    """从代码内容推断文件路径"""
    if lang == 'java':
        # 尝试从 package 和 class 声明推断路径
        package_match = re.search(r'package\s+([\w.]+);', code)
        class_match = re.search(r'(?:public\s+)?(?:class|interface|enum)\s+(\w+)', code)
        
        if class_match:
            class_name = class_match.group(1)
            if package_match:
                package_path = package_match.group(1).replace('.', '/')
                return f'src/main/java/{package_path}/{class_name}.java'
            else:
                return f'{class_name}.java'
    
    return None


def _save_code_project(workspace_id: str, content: str, project_name: str = 'project') -> list[str]:
    """从 AI 输出中提取代码块并保存为项目文件"""
    workspace_path = WORKSPACES_ROOT / workspace_id
    saved_files = []
    
    code_blocks = _extract_code_blocks(content)
    
    if not code_blocks:
        # 没有找到代码块，保存为 markdown
        return []
    
    # 语言到扩展名的映射
    lang_ext_map = {
        'java': '.java',
        'python': '.py', 'py': '.py',
        'javascript': '.js', 'js': '.js',
        'typescript': '.ts', 'ts': '.ts',
        'go': '.go',
        'rust': '.rs', 'rs': '.rs',
        'cpp': '.cpp', 'c++': '.cpp',
        'c': '.c',
        'sql': '.sql',
        'xml': '.xml',
        'yaml': '.yaml', 'yml': '.yaml',
        'json': '.json',
        'properties': '.properties',
        'html': '.html',
        'css': '.css',
        'sh': '.sh', 'bash': '.sh', 'shell': '.sh',
    }
    
    for i, block in enumerate(code_blocks):
        lang = block['language'].lower()
        filepath = block['filepath']
        code = block['code']
        
        if filepath:
            # 有指定路径，按路径保存
            target_path = workspace_path / 'code' / project_name / filepath
        else:
            # 没有指定路径，根据语言生成文件名
            ext = lang_ext_map.get(lang, f'.{lang}')
            target_path = workspace_path / 'code' / project_name / f'file_{i+1}{ext}'
        
        target_path.parent.mkdir(parents=True, exist_ok=True)
        target_path.write_text(code, encoding='utf-8')
        saved_files.append(str(target_path.relative_to(workspace_path)))
    
    return saved_files


def _generate_filename(message: str, ext: str) -> str:
    """根据消息生成文件名"""
    clean = re.sub(r'^(帮我|请|生成|写|创建|制作|编写|做一个?|给我)', '', message)
    name = clean[:20].strip()
    name = re.sub(r'[\\/:*?"<>|\s]', '_', name)
    timestamp = datetime.now().strftime("%H%M%S")
    
    if not name:
        name = f"output_{timestamp}"
    else:
        name = f"{name}_{timestamp}"
    
    return f"{name}.{ext}"


def _save_markdown(workspace_id: str, content: str, folder: str, filename: str) -> str:
    """保存 Markdown/文本内容"""
    workspace_path = WORKSPACES_ROOT / workspace_id
    target_dir = workspace_path / folder
    target_dir.mkdir(parents=True, exist_ok=True)
    
    file_path = target_dir / filename
    file_path.write_text(content, encoding='utf-8')
    
    return f"{folder}/{filename}"


async def _generate_slide_image(title: str, workspace_path: Path, slide_index: int) -> Optional[Path]:
    """为幻灯片生成配图"""
    try:
        from auto.core.ai.image import get_image_generator
        
        generator = get_image_generator()
        
        # 生成图像描述
        prompt = f"Professional business presentation slide background image for topic: {title}. Modern, clean, minimalist design with subtle gradient. No text. Corporate style."
        
        image_path = workspace_path / "images" / f"slide_{slide_index}.png"
        result = await generator.generate_and_download(prompt, image_path, size="1792x1024")
        
        return result
    except Exception as e:
        print(f"图像生成跳过: {e}")
        return None


def _generate_pptx_sync(workspace_id: str, content: str, filename: str, with_images: bool = False) -> str:
    """同步生成 PPTX 文件（不含图片）"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        # 解析内容生成幻灯片
        lines = content.split('\n')
        current_title = ""
        current_content = []
        slides_data = []
        
        for line in lines:
            line = line.strip()
            if line.startswith('# '):
                if current_title or current_content:
                    slides_data.append((current_title, current_content))
                current_title = line[2:]
                current_content = []
            elif line.startswith('## '):
                if current_title or current_content:
                    slides_data.append((current_title, current_content))
                current_title = line[3:]
                current_content = []
            elif line.startswith('### '):
                current_content.append(f"◆ {line[4:]}")
            elif line.startswith('- ') or line.startswith('* '):
                current_content.append(f"• {line[2:]}")
            elif line.startswith('|') or line.startswith('```'):
                continue
            elif line:
                current_content.append(line)
        
        if current_title or current_content:
            slides_data.append((current_title, current_content))
        
        # 颜色方案
        primary_color = RGBColor(0x1E, 0x3A, 0x8A)  # 深蓝色
        accent_color = RGBColor(0x3B, 0x82, 0xF6)   # 亮蓝色
        
        # 生成幻灯片
        for i, (title, content_lines) in enumerate(slides_data):
            slide_layout = prs.slide_layouts[6]  # 空白布局
            slide = prs.slides.add_slide(slide_layout)
            
            # 添加背景色块（左侧装饰条）
            left_bar = slide.shapes.add_shape(
                1,  # MSO_SHAPE.RECTANGLE
                Inches(0), Inches(0),
                Inches(0.3), Inches(7.5)
            )
            left_bar.fill.solid()
            left_bar.fill.fore_color.rgb = primary_color
            left_bar.line.fill.background()
            
            if i == 0:
                # 封面页
                title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.5), Inches(1.5))
                tf = title_box.text_frame
                p = tf.paragraphs[0]
                p.text = title
                p.font.size = Pt(48)
                p.font.bold = True
                p.font.color.rgb = primary_color
                p.alignment = PP_ALIGN.CENTER
                
                # 副标题
                if content_lines:
                    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(11.5), Inches(1))
                    tf = sub_box.text_frame
                    p = tf.paragraphs[0]
                    p.text = content_lines[0] if content_lines else ""
                    p.font.size = Pt(24)
                    p.font.color.rgb = accent_color
                    p.alignment = PP_ALIGN.CENTER
            else:
                # 内容页标题
                title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.8))
                tf = title_box.text_frame
                p = tf.paragraphs[0]
                p.text = title
                p.font.size = Pt(32)
                p.font.bold = True
                p.font.color.rgb = primary_color
                
                # 分隔线
                line = slide.shapes.add_shape(
                    1, Inches(0.8), Inches(1.15), Inches(11.5), Inches(0.03)
                )
                line.fill.solid()
                line.fill.fore_color.rgb = accent_color
                line.line.fill.background()
                
                # 内容
                if content_lines:
                    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.5), Inches(5.5))
                    tf = content_box.text_frame
                    tf.word_wrap = True
                    
                    for j, line_text in enumerate(content_lines[:10]):
                        if j == 0:
                            p = tf.paragraphs[0]
                        else:
                            p = tf.add_paragraph()
                        p.text = line_text
                        p.font.size = Pt(20)
                        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                        p.space_after = Pt(12)
        
        # 保存
        workspace_path = WORKSPACES_ROOT / workspace_id
        target_dir = workspace_path / "ppt"
        target_dir.mkdir(parents=True, exist_ok=True)
        
        file_path = target_dir / filename
        prs.save(str(file_path))
        
        return f"ppt/{filename}"
    except Exception as e:
        import traceback
        print(f"PPT生成失败: {e}")
        traceback.print_exc()
        return _save_markdown(workspace_id, content, "docs", filename.replace('.pptx', '.md'))


def _generate_pptx(workspace_id: str, content: str, filename: str) -> str:
    """生成 PPTX 文件（同步版本）"""
    return _generate_pptx_sync(workspace_id, content, filename)


def _generate_xlsx(workspace_id: str, content: str, filename: str) -> str:
    """生成真正的 XLSX 文件"""
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
        
        wb = Workbook()
        ws = wb.active
        ws.title = "数据"
        
        # 解析 Markdown 表格或内容
        lines = content.split('\n')
        row = 1
        
        # 样式
        header_font = Font(bold=True, color="FFFFFF")
        header_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
        border = Border(
            left=Side(style='thin'),
            right=Side(style='thin'),
            top=Side(style='thin'),
            bottom=Side(style='thin')
        )
        
        for line in lines:
            line = line.strip()
            if line.startswith('|') and '|' in line[1:]:
                # Markdown 表格行
                if '---' in line:
                    continue
                cells = [c.strip() for c in line.split('|')[1:-1]]
                for col, value in enumerate(cells, 1):
                    cell = ws.cell(row=row, column=col, value=value)
                    cell.border = border
                    if row == 1:
                        cell.font = header_font
                        cell.fill = header_fill
                    cell.alignment = Alignment(horizontal='center', vertical='center')
                row += 1
            elif line and not line.startswith('#') and not line.startswith('```'):
                # 普通内容
                ws.cell(row=row, column=1, value=line)
                row += 1
        
        # 调整列宽
        for col in ws.columns:
            max_length = 0
            column = col[0].column_letter
            for cell in col:
                try:
                    if len(str(cell.value)) > max_length:
                        max_length = len(str(cell.value))
                except:
                    pass
            ws.column_dimensions[column].width = min(max_length + 2, 50)
        
        # 保存
        workspace_path = WORKSPACES_ROOT / workspace_id
        target_dir = workspace_path / "data"
        target_dir.mkdir(parents=True, exist_ok=True)
        
        file_path = target_dir / filename
        wb.save(str(file_path))
        
        return f"data/{filename}"
    except ImportError:
        return _save_markdown(workspace_id, content, "docs", filename.replace('.xlsx', '.md'))


@router.post("/chat")
async def chat(request: ChatRequest) -> ChatResponse:
    """发送聊天消息"""
    try:
        router_instance = get_router()
        
        # 提取核心消息（去掉参考文件标记等）
        core_message = request.message
        if '[参考文件' in core_message:
            # 去掉参考文件标记，只保留实际请求
            import re
            core_message = re.sub(r'\[参考文件[^\]]*\]\s*', '', core_message).strip()
        
        # 检测文件类型和是否需要保存
        ext, folder = _detect_file_type(core_message)
        should_save = _should_save_file(core_message)
        
        # ========== 读取工作空间知识上下文 ==========
        workspace_context = ""
        if request.workspace_id:
            try:
                from auto.core.workspace.prompts import get_workspace_prompts
                
                prompts = get_workspace_prompts(request.workspace_id)
                if prompts:
                    # 检查是否需要读取项目文件（关键词检测）
                    need_context = any(kw in core_message.lower() for kw in [
                        '项目', '代码', '文件', '分析', '学习', '理解', '阅读', '读取',
                        'project', 'code', 'file', 'analyze', 'learn', 'understand', 'read',
                        '帮我看', '看看', '检查', '优化', '重构', '解释',
                    ])
                    
                    if need_context:
                        # 读取工作空间所有文件作为上下文
                        workspace_context = prompts.build_knowledge_context(
                            include_file_contents=True,
                            max_files=30,
                        )
                    else:
                        # 只读取文件结构摘要
                        workspace_context = prompts.build_knowledge_context(
                            include_file_contents=False,
                            summary_only=True,
                        )
            except Exception as e:
                print(f"读取工作空间上下文失败: {e}")
        
        # 根据文件类型调整 AI 提示（仅当需要保存文件时）
        if should_save and ext == 'pptx':
            prompt = f"""你是一个专业的PPT设计师。请为以下主题创建完整的PPT内容。

要求：
1. 严格使用以下Markdown格式输出
2. 用 # 表示PPT标题（封面）
3. 用 ## 表示每一页的标题
4. 用 - 表示每页的要点内容
5. 生成 8-12 页内容
6. 内容要专业、具体、有数据支撑

格式示例：
# 产品发布会

## 公司简介
- 成立于2020年
- 专注于xxx领域
- 服务超过100万用户

## 产品亮点
- 功能一：描述
- 功能二：描述

主题：{core_message}

请直接输出PPT内容，不要有任何解释性文字："""
        elif should_save and ext == 'xlsx':
            prompt = f"""请为以下需求创建完整的表格数据。使用 Markdown 表格格式，包含表头和至少10行数据：

| 列1 | 列2 | 列3 | 列4 |
|-----|-----|-----|-----|
| 数据 | 数据 | 数据 | 数据 |

需求：{core_message}

请直接输出表格，不要有任何解释性文字："""
        elif should_save and ext == 'sql':
            prompt = f"""请为以下需求编写 SQL 语句。

要求：
1. 使用标准 SQL 语法
2. 添加必要的注释说明
3. 考虑数据类型和约束
4. 如果是建表语句，包含主键、索引等

需求：{core_message}

请直接输出 SQL 代码："""
        elif should_save and ext in ('java_project', 'code_project'):
            prompt = f"""请根据工作空间中的产品文档和数据库设计，生成完整的后端代码。

重要：请按照以下格式输出每个代码文件，确保每个代码块都指定文件路径：

```java:src/main/java/com/rental/entity/User.java
package com.rental.entity;

// 代码内容
```

```java:src/main/java/com/rental/controller/UserController.java
package com.rental.controller;

// 代码内容
```

要求：
1. 每个代码块必须使用 ```language:filepath 格式
2. 生成完整的、可运行的代码
3. 包含实体类、Controller、Service、Mapper 等完整结构
4. 添加必要的注解和注释
5. 符合 Java/Spring 最佳实践

需求：{core_message}

请生成完整代码："""
        else:
            prompt = request.message
        
        # 构建消息列表
        messages = []
        
        # 如果有工作空间上下文，先添加系统消息
        if workspace_context:
            messages.append(Message(
                role=MessageRole.SYSTEM,
                content=f"""你已经阅读并学习了当前工作空间中的所有文件。请根据这些知识来回答用户的问题。

{workspace_context}

---

请根据以上工作空间的内容来理解项目，并帮助用户完成任务。"""
            ))
        
        messages.append(Message(role=MessageRole.USER, content=prompt))
        
        response = await router_instance.chat(
            messages=messages,
            model=request.model,
        )
        
        content = response.message.content
        saved_file = None
        
        # 判断是否需要保存文件
        should_save = _should_save_file(core_message)
        
        # 如果有工作空间 ID 且需要保存且用户允许保存
        if request.workspace_id and request.save_to_workspace and should_save:
            workspace_path = WORKSPACES_ROOT / request.workspace_id
            if workspace_path.exists():
                filename = _generate_filename(core_message, ext)
                
                # 根据类型生成不同格式的文件
                if ext == 'pptx':
                    saved_file = _generate_pptx(request.workspace_id, content, filename)
                elif ext == 'xlsx':
                    saved_file = _generate_xlsx(request.workspace_id, content, filename)
                elif ext in ('java_project', 'code_project'):
                    # 代码项目：提取所有代码块并分别保存
                    project_name = re.sub(r'[\\/:*?"<>|\s]', '_', core_message[:20].strip()) or 'project'
                    saved_files = _save_code_project(request.workspace_id, content, project_name)
                    if saved_files:
                        saved_file = ', '.join(saved_files[:5])  # 最多显示5个文件
                        if len(saved_files) > 5:
                            saved_file += f' 等 {len(saved_files)} 个文件'
                    else:
                        # 没有提取到代码块，保存为 markdown
                        saved_file = _save_markdown(request.workspace_id, content, 'docs', filename.replace(f'.{ext}', '.md'))
                elif ext == 'sql':
                    saved_file = _save_markdown(request.workspace_id, content, folder, filename)
                elif ext in ('py', 'js', 'ts', 'java', 'go', 'rs', 'cpp', 'sh', 'html'):
                    saved_file = _save_markdown(request.workspace_id, content, folder, filename)
                elif ext in ('json', 'yaml', 'xml'):
                    saved_file = _save_markdown(request.workspace_id, content, folder, filename)
                else:
                    saved_file = _save_markdown(request.workspace_id, content, folder, filename)
        
        return ChatResponse(
            data={
                "id": response.id,
                "content": content,
                "model": response.model,
                "saved_file": saved_file,
                "file_type": ext,
                "usage": {
                    "prompt_tokens": response.usage.prompt_tokens,
                    "completion_tokens": response.usage.completion_tokens,
                    "total_tokens": response.usage.total_tokens,
                },
            }
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/conversations/{conversation_id}/messages")
async def send_message(
    conversation_id: str,
    request: ChatRequest,
) -> ChatResponse:
    """在会话中发送消息"""
    # TODO: 实现会话管理
    return await chat(request)
