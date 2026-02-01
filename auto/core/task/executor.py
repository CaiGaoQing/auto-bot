"""任务执行器 - 生成真实的 PPT/Excel/PDF 文件"""

import os
import json
from datetime import datetime
from pathlib import Path
from typing import Optional, List, Dict, Any
from dataclasses import dataclass
from enum import Enum


class OutputFormat(str, Enum):
    """输出格式"""
    MARKDOWN = "md"
    EXCEL = "xlsx"
    PPT = "pptx"
    PDF = "pdf"
    JSON = "json"


@dataclass
class TaskResult:
    """任务执行结果"""
    success: bool
    message: str
    output_files: List[Dict[str, Any]]
    content: Optional[str] = None
    error: Optional[str] = None


class TaskExecutor:
    """任务执行器 - 调用 AI 并生成真实文件"""
    
    def __init__(self, workspace_path: Path):
        self.workspace_path = workspace_path
        self.files_path = workspace_path / "files"
        self.outputs_path = workspace_path / "outputs"
        self.files_path.mkdir(parents=True, exist_ok=True)
        self.outputs_path.mkdir(parents=True, exist_ok=True)
    
    async def execute(
        self,
        task: str,
        input_files: List[str] = [],
        output_format: OutputFormat = OutputFormat.MARKDOWN,
    ) -> TaskResult:
        """执行任务并生成输出文件"""
        try:
            # 1. 读取输入文件内容
            file_contents = self._read_input_files(input_files)
            
            # 2. 调用 AI 分析
            ai_result = await self._call_ai(task, file_contents, output_format)
            
            if not ai_result:
                return TaskResult(
                    success=False,
                    message="AI 处理失败",
                    output_files=[],
                    error="无法获取 AI 响应"
                )
            
            # 3. 生成输出文件
            output_files = await self._generate_output(task, ai_result, output_format)
            
            return TaskResult(
                success=True,
                message="任务执行完成",
                output_files=output_files,
                content=ai_result[:2000] if ai_result else None,
            )
            
        except Exception as e:
            return TaskResult(
                success=False,
                message="执行失败",
                output_files=[],
                error=str(e)
            )
    
    def _read_input_files(self, filenames: List[str]) -> Dict[str, str]:
        """读取输入文件内容"""
        contents = {}
        
        for filename in filenames:
            file_path = self.files_path / filename
            if not file_path.exists():
                continue
            
            suffix = file_path.suffix.lower()
            
            try:
                if suffix in ['.txt', '.md', '.csv', '.json', '.log']:
                    contents[filename] = file_path.read_text(encoding='utf-8')[:50000]
                
                elif suffix in ['.xlsx', '.xls']:
                    try:
                        import pandas as pd
                        df = pd.read_excel(file_path)
                        contents[filename] = f"[Excel 文件: {filename}]\n列: {list(df.columns)}\n行数: {len(df)}\n\n数据预览:\n{df.to_string()}"
                    except ImportError:
                        contents[filename] = f"[Excel 文件: {filename}]"
                
                elif suffix == '.pdf':
                    try:
                        import pymupdf
                        doc = pymupdf.open(file_path)
                        text = ""
                        for page in doc[:20]:  # 最多20页
                            text += page.get_text()
                        contents[filename] = f"[PDF 文件: {filename}]\n\n{text[:50000]}"
                    except ImportError:
                        contents[filename] = f"[PDF 文件: {filename}]"
                
                elif suffix in ['.docx']:
                    try:
                        from docx import Document
                        doc = Document(file_path)
                        text = "\n".join([p.text for p in doc.paragraphs])
                        contents[filename] = f"[Word 文件: {filename}]\n\n{text[:50000]}"
                    except ImportError:
                        contents[filename] = f"[Word 文件: {filename}]"
                
                else:
                    contents[filename] = f"[文件: {filename}, 类型: {suffix}]"
                    
            except Exception as e:
                contents[filename] = f"[读取文件失败: {filename}, 错误: {str(e)}]"
        
        return contents
    
    async def _call_ai(
        self,
        task: str,
        file_contents: Dict[str, str],
        output_format: OutputFormat,
    ) -> Optional[str]:
        """调用 AI 处理任务"""
        from auto.core.ai.router import get_router
        from auto.shared.models import Message, MessageRole
        
        # 构建文件上下文
        file_context = ""
        if file_contents:
            file_context = "\n\n## 输入文件内容:\n"
            for name, content in file_contents.items():
                # 限制每个文件的内容长度
                truncated = content[:20000] if len(content) > 20000 else content
                file_context += f"\n### 文件: {name}\n```\n{truncated}\n```\n"
        
        # 根据输出格式调整提示
        format_instructions = {
            OutputFormat.MARKDOWN: "请以 Markdown 格式输出结果。",
            OutputFormat.EXCEL: """请以结构化表格格式输出结果。使用以下 JSON 格式:
```json
{
  "sheets": [
    {
      "name": "Sheet1",
      "headers": ["列1", "列2", "列3"],
      "data": [
        ["值1", "值2", "值3"],
        ["值4", "值5", "值6"]
      ]
    }
  ]
}
```""",
            OutputFormat.PPT: """请以演示文稿格式输出结果。使用以下 JSON 格式:
```json
{
  "title": "演示标题",
  "slides": [
    {
      "title": "幻灯片标题",
      "content": ["要点1", "要点2", "要点3"]
    }
  ]
}
```""",
            OutputFormat.PDF: "请以专业报告格式输出结果，使用 Markdown 格式，包含标题、章节、列表等。",
            OutputFormat.JSON: "请以 JSON 格式输出结构化结果。",
        }
        
        prompt = f"""你是一个专业的工作助手。请完成以下任务：

## 任务描述:
{task}
{file_context}

## 输出要求:
{format_instructions.get(output_format, '')}

请详细分析并给出专业的结果。
"""
        
        try:
            router = get_router()
            response = await router.chat(
                messages=[Message(role=MessageRole.USER, content=prompt)],
            )
            return response.message.content
        except Exception as e:
            print(f"AI 调用失败: {e}")
            return None
    
    async def _generate_output(
        self,
        task: str,
        ai_result: str,
        output_format: OutputFormat,
    ) -> List[Dict[str, Any]]:
        """根据 AI 结果生成输出文件"""
        output_files = []
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        
        if output_format == OutputFormat.EXCEL:
            file_info = self._generate_excel(task, ai_result, timestamp)
            if file_info:
                output_files.append(file_info)
        
        elif output_format == OutputFormat.PPT:
            file_info = self._generate_ppt(task, ai_result, timestamp)
            if file_info:
                output_files.append(file_info)
        
        elif output_format == OutputFormat.PDF:
            file_info = self._generate_pdf(task, ai_result, timestamp)
            if file_info:
                output_files.append(file_info)
        
        elif output_format == OutputFormat.JSON:
            file_info = self._generate_json(task, ai_result, timestamp)
            if file_info:
                output_files.append(file_info)
        
        # 始终生成 Markdown 备份
        md_info = self._generate_markdown(task, ai_result, timestamp)
        if md_info:
            output_files.append(md_info)
        
        return output_files
    
    def _generate_excel(self, task: str, ai_result: str, timestamp: str) -> Optional[Dict]:
        """生成 Excel 文件"""
        try:
            import pandas as pd
            from openpyxl import Workbook
            from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
            
            output_file = self.outputs_path / f"result_{timestamp}.xlsx"
            
            # 尝试解析 JSON 格式的表格数据
            json_data = self._extract_json(ai_result)
            
            if json_data and 'sheets' in json_data:
                # AI 返回了结构化数据
                wb = Workbook()
                wb.remove(wb.active)
                
                for sheet_data in json_data['sheets']:
                    ws = wb.create_sheet(title=sheet_data.get('name', 'Sheet')[:31])
                    
                    # 写入表头
                    headers = sheet_data.get('headers', [])
                    for col, header in enumerate(headers, 1):
                        cell = ws.cell(row=1, column=col, value=str(header))
                        cell.font = Font(bold=True, color="FFFFFF")
                        cell.fill = PatternFill(start_color="2563EB", end_color="2563EB", fill_type="solid")
                        cell.alignment = Alignment(horizontal='center')
                    
                    # 写入数据
                    for row_idx, row_data in enumerate(sheet_data.get('data', []), 2):
                        for col_idx, value in enumerate(row_data, 1):
                            ws.cell(row=row_idx, column=col_idx, value=value)
                    
                    # 调整列宽
                    for col in ws.columns:
                        max_length = 0
                        for cell in col:
                            try:
                                if len(str(cell.value)) > max_length:
                                    max_length = len(str(cell.value))
                            except:
                                pass
                        ws.column_dimensions[col[0].column_letter].width = min(max_length + 2, 50)
                
                wb.save(output_file)
            else:
                # 将 AI 结果作为文本放入 Excel
                wb = Workbook()
                ws = wb.active
                ws.title = "分析结果"
                
                # 标题
                ws['A1'] = task[:100]
                ws['A1'].font = Font(bold=True, size=14)
                
                # 内容
                lines = ai_result.split('\n')
                for i, line in enumerate(lines[:1000], 3):
                    ws.cell(row=i, column=1, value=line)
                
                ws.column_dimensions['A'].width = 100
                wb.save(output_file)
            
            return {
                "name": output_file.name,
                "path": str(output_file),
                "size": output_file.stat().st_size,
                "type": "xlsx",
            }
            
        except Exception as e:
            print(f"生成 Excel 失败: {e}")
            return None
    
    def _generate_ppt(self, task: str, ai_result: str, timestamp: str) -> Optional[Dict]:
        """生成 PPT 文件"""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RgbColor
            from pptx.enum.text import PP_ALIGN
            
            output_file = self.outputs_path / f"result_{timestamp}.pptx"
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            
            # 尝试解析 JSON 格式的 PPT 数据
            json_data = self._extract_json(ai_result)
            
            if json_data and 'slides' in json_data:
                # AI 返回了结构化数据
                ppt_title = json_data.get('title', task[:50])
                slides_data = json_data.get('slides', [])
                
                # 标题页
                slide_layout = prs.slide_layouts[6]  # 空白布局
                slide = prs.slides.add_slide(slide_layout)
                
                # 添加标题
                title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
                tf = title_box.text_frame
                p = tf.paragraphs[0]
                p.text = ppt_title
                p.font.size = Pt(44)
                p.font.bold = True
                p.alignment = PP_ALIGN.CENTER
                
                # 添加日期
                date_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.333), Inches(0.5))
                tf = date_box.text_frame
                p = tf.paragraphs[0]
                p.text = datetime.now().strftime("%Y年%m月%d日")
                p.font.size = Pt(20)
                p.alignment = PP_ALIGN.CENTER
                
                # 内容页
                for slide_data in slides_data[:20]:  # 最多20页
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    
                    # 幻灯片标题
                    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
                    tf = title_box.text_frame
                    p = tf.paragraphs[0]
                    p.text = slide_data.get('title', '')
                    p.font.size = Pt(32)
                    p.font.bold = True
                    
                    # 内容
                    content = slide_data.get('content', [])
                    if isinstance(content, list):
                        content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(11.933), Inches(5.5))
                        tf = content_box.text_frame
                        tf.word_wrap = True
                        
                        for i, item in enumerate(content):
                            if i == 0:
                                p = tf.paragraphs[0]
                            else:
                                p = tf.add_paragraph()
                            p.text = f"• {item}"
                            p.font.size = Pt(24)
                            p.space_before = Pt(12)
                    else:
                        content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(11.933), Inches(5.5))
                        tf = content_box.text_frame
                        tf.word_wrap = True
                        p = tf.paragraphs[0]
                        p.text = str(content)
                        p.font.size = Pt(20)
            else:
                # 从文本自动生成 PPT
                # 标题页
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
                tf = title_box.text_frame
                p = tf.paragraphs[0]
                p.text = task[:60]
                p.font.size = Pt(40)
                p.font.bold = True
                p.alignment = PP_ALIGN.CENTER
                
                # 按段落分割内容
                paragraphs = [p.strip() for p in ai_result.split('\n\n') if p.strip()]
                
                current_slide = None
                current_content = []
                
                for para in paragraphs[:50]:
                    # 检测是否是标题（以 # 开头或较短的行）
                    is_title = para.startswith('#') or (len(para) < 50 and not para.startswith('-') and not para.startswith('•'))
                    
                    if is_title and len(para) < 100:
                        # 保存之前的幻灯片
                        if current_slide and current_content:
                            self._add_ppt_content(current_slide, current_content)
                        
                        # 创建新幻灯片
                        current_slide = prs.slides.add_slide(prs.slide_layouts[6])
                        title_text = para.lstrip('#').strip()
                        
                        title_box = current_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
                        tf = title_box.text_frame
                        p = tf.paragraphs[0]
                        p.text = title_text[:80]
                        p.font.size = Pt(32)
                        p.font.bold = True
                        
                        current_content = []
                    else:
                        current_content.append(para)
                
                # 保存最后一个幻灯片
                if current_slide and current_content:
                    self._add_ppt_content(current_slide, current_content)
            
            prs.save(output_file)
            
            return {
                "name": output_file.name,
                "path": str(output_file),
                "size": output_file.stat().st_size,
                "type": "pptx",
            }
            
        except Exception as e:
            print(f"生成 PPT 失败: {e}")
            import traceback
            traceback.print_exc()
            return None
    
    def _add_ppt_content(self, slide, content_list):
        """向 PPT 幻灯片添加内容"""
        from pptx.util import Inches, Pt
        
        content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(11.933), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, content in enumerate(content_list[:10]):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            # 处理列表项
            if content.startswith('-') or content.startswith('•'):
                p.text = content
            else:
                p.text = content[:500]
            
            p.font.size = Pt(18)
            p.space_before = Pt(8)
    
    def _generate_pdf(self, task: str, ai_result: str, timestamp: str) -> Optional[Dict]:
        """生成 PDF 文件"""
        try:
            from weasyprint import HTML, CSS
            import markdown
            
            output_file = self.outputs_path / f"result_{timestamp}.pdf"
            
            # 将 Markdown 转换为 HTML
            html_content = markdown.markdown(
                ai_result,
                extensions=['tables', 'fenced_code', 'codehilite']
            )
            
            # 创建完整的 HTML 文档
            html = f"""
            <!DOCTYPE html>
            <html>
            <head>
                <meta charset="utf-8">
                <title>{task[:50]}</title>
            </head>
            <body>
                <h1 style="color: #2563EB; border-bottom: 2px solid #2563EB; padding-bottom: 10px;">
                    {task[:100]}
                </h1>
                <p style="color: #666; font-size: 12px; margin-bottom: 30px;">
                    生成时间: {datetime.now().strftime("%Y-%m-%d %H:%M:%S")}
                </p>
                {html_content}
            </body>
            </html>
            """
            
            # PDF 样式
            css = CSS(string="""
                @page {
                    size: A4;
                    margin: 2cm;
                }
                body {
                    font-family: "Noto Sans SC", "Microsoft YaHei", sans-serif;
                    font-size: 12pt;
                    line-height: 1.6;
                    color: #333;
                }
                h1 { font-size: 24pt; margin-top: 0; }
                h2 { font-size: 18pt; color: #2563EB; margin-top: 20pt; }
                h3 { font-size: 14pt; color: #3B82F6; margin-top: 15pt; }
                p { margin: 10pt 0; }
                ul, ol { margin: 10pt 0; padding-left: 20pt; }
                li { margin: 5pt 0; }
                table {
                    width: 100%;
                    border-collapse: collapse;
                    margin: 15pt 0;
                }
                th, td {
                    border: 1px solid #ddd;
                    padding: 8pt;
                    text-align: left;
                }
                th {
                    background-color: #2563EB;
                    color: white;
                }
                tr:nth-child(even) { background-color: #f9f9f9; }
                code {
                    background-color: #f5f5f5;
                    padding: 2pt 4pt;
                    border-radius: 3pt;
                    font-family: monospace;
                }
                pre {
                    background-color: #f5f5f5;
                    padding: 10pt;
                    border-radius: 5pt;
                    overflow-x: auto;
                }
            """)
            
            HTML(string=html).write_pdf(output_file, stylesheets=[css])
            
            return {
                "name": output_file.name,
                "path": str(output_file),
                "size": output_file.stat().st_size,
                "type": "pdf",
            }
            
        except Exception as e:
            print(f"生成 PDF 失败: {e}")
            return None
    
    def _generate_json(self, task: str, ai_result: str, timestamp: str) -> Optional[Dict]:
        """生成 JSON 文件"""
        try:
            output_file = self.outputs_path / f"result_{timestamp}.json"
            
            # 尝试从 AI 结果中提取 JSON
            json_data = self._extract_json(ai_result)
            
            if json_data:
                result = json_data
            else:
                result = {
                    "task": task,
                    "result": ai_result,
                    "generated_at": datetime.now().isoformat(),
                }
            
            output_file.write_text(
                json.dumps(result, ensure_ascii=False, indent=2),
                encoding='utf-8'
            )
            
            return {
                "name": output_file.name,
                "path": str(output_file),
                "size": output_file.stat().st_size,
                "type": "json",
            }
            
        except Exception as e:
            print(f"生成 JSON 失败: {e}")
            return None
    
    def _generate_markdown(self, task: str, ai_result: str, timestamp: str) -> Optional[Dict]:
        """生成 Markdown 文件"""
        try:
            output_file = self.outputs_path / f"result_{timestamp}.md"
            
            content = f"""# {task}

> 生成时间: {datetime.now().strftime("%Y-%m-%d %H:%M:%S")}

---

{ai_result}
"""
            output_file.write_text(content, encoding='utf-8')
            
            return {
                "name": output_file.name,
                "path": str(output_file),
                "size": output_file.stat().st_size,
                "type": "md",
            }
            
        except Exception as e:
            print(f"生成 Markdown 失败: {e}")
            return None
    
    def _extract_json(self, text: str) -> Optional[Dict]:
        """从文本中提取 JSON"""
        import re
        
        # 尝试找到 JSON 代码块
        json_pattern = r'```json\s*([\s\S]*?)\s*```'
        matches = re.findall(json_pattern, text)
        
        for match in matches:
            try:
                return json.loads(match)
            except:
                continue
        
        # 尝试直接解析整个文本
        try:
            return json.loads(text)
        except:
            pass
        
        return None
