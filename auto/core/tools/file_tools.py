"""
文件操作工具

提供给 AI 调用的文件创建和操作工具
"""

from pathlib import Path
from typing import Optional, List
from datetime import datetime
import json

from .base import BaseTool, ToolResult, ToolParameter


# 工作空间根目录
_PROJECT_ROOT = Path(__file__).parent.parent.parent.parent
WORKSPACES_ROOT = _PROJECT_ROOT / "data" / "workspaces"


class CreateFileTool(BaseTool):
    """
    创建文件工具
    
    让 AI 能够在工作空间中创建任意类型的文件
    """
    
    name = "create_file"
    display_name = "创建文件"
    description = """在工作空间中创建文件。
    
使用场景：
- 创建代码文件（.java, .py, .ts 等）
- 创建配置文件（.json, .yaml 等）
- 创建文档文件（.md, .txt 等）
- 创建 SQL 脚本（.sql）

参数说明：
- filepath: 文件路径，相对于工作空间根目录，如 "src/main/java/com/example/User.java"
- content: 文件内容
- folder: 可选的顶级文件夹，如 "code", "scripts", "docs"
"""
    category = "file"
    
    parameters = [
        ToolParameter(
            name="filepath",
            description="文件路径（相对于工作空间），如 'src/User.java' 或 'scripts/init.sql'",
            type="string",
            required=True,
        ),
        ToolParameter(
            name="content",
            description="文件内容",
            type="string",
            required=True,
        ),
        ToolParameter(
            name="folder",
            description="顶级文件夹（code/scripts/docs/data），默认根据文件类型自动选择",
            type="string",
            required=False,
            enum=["code", "scripts", "docs", "data", "ppt", ""],
        ),
    ]
    
    async def execute(
        self,
        filepath: str,
        content: str,
        folder: Optional[str] = None,
    ) -> ToolResult:
        """创建文件"""
        if not self.workspace_id:
            return ToolResult(
                success=False,
                message="未指定工作空间",
                error="workspace_id is required",
            )
        
        workspace_path = WORKSPACES_ROOT / self.workspace_id
        if not workspace_path.exists():
            return ToolResult(
                success=False,
                message=f"工作空间不存在: {self.workspace_id}",
                error="workspace not found",
            )
        
        # 自动选择文件夹
        if not folder:
            folder = self._detect_folder(filepath)
        
        # 构建完整路径
        if folder:
            full_path = workspace_path / folder / filepath
        else:
            full_path = workspace_path / filepath
        
        # 创建目录
        full_path.parent.mkdir(parents=True, exist_ok=True)
        
        # 写入文件
        try:
            full_path.write_text(content, encoding="utf-8")
            
            relative_path = str(full_path.relative_to(workspace_path))
            
            return ToolResult(
                success=True,
                message=f"文件已创建: {relative_path}",
                files_created=[relative_path],
                data={
                    "filepath": relative_path,
                    "size": len(content),
                }
            )
        except Exception as e:
            return ToolResult(
                success=False,
                message="文件创建失败",
                error=str(e),
            )
    
    def _detect_folder(self, filepath: str) -> str:
        """根据文件类型检测文件夹"""
        ext = Path(filepath).suffix.lower()
        
        code_exts = {'.java', '.py', '.js', '.ts', '.tsx', '.go', '.rs', '.cpp', '.c', '.h', '.html', '.css', '.vue', '.svelte'}
        script_exts = {'.sql', '.sh', '.bash', '.ps1', '.bat'}
        doc_exts = {'.md', '.txt', '.rst', '.adoc'}
        data_exts = {'.json', '.yaml', '.yml', '.xml', '.csv', '.properties'}
        
        if ext in code_exts:
            return "code"
        elif ext in script_exts:
            return "scripts"
        elif ext in doc_exts:
            return "docs"
        elif ext in data_exts:
            return "data"
        else:
            return "code"  # 默认放代码目录


class SaveCodeTool(BaseTool):
    """
    保存代码项目工具
    
    批量创建多个代码文件，适合生成完整项目结构
    """
    
    name = "save_code_project"
    display_name = "保存代码项目"
    description = """批量创建多个代码文件，用于生成完整的项目结构。
    
使用场景：
- 生成 SpringBoot 项目代码
- 生成 React/Vue 项目代码
- 生成微服务项目
- 任何需要多个文件的代码项目

参数说明：
- project_name: 项目名称，将作为顶级目录
- files: 文件列表，每个文件包含 path 和 content
"""
    category = "file"
    
    parameters = [
        ToolParameter(
            name="project_name",
            description="项目名称",
            type="string",
            required=True,
        ),
        ToolParameter(
            name="files",
            description="文件列表，格式: [{\"path\": \"src/User.java\", \"content\": \"...\"}]",
            type="array",
            required=True,
        ),
    ]
    
    async def execute(
        self,
        project_name: str,
        files: List[dict],
    ) -> ToolResult:
        """保存代码项目"""
        if not self.workspace_id:
            return ToolResult(
                success=False,
                message="未指定工作空间",
                error="workspace_id is required",
            )
        
        workspace_path = WORKSPACES_ROOT / self.workspace_id
        if not workspace_path.exists():
            return ToolResult(
                success=False,
                message=f"工作空间不存在: {self.workspace_id}",
                error="workspace not found",
            )
        
        # 项目根目录
        project_path = workspace_path / "code" / project_name
        
        created_files = []
        errors = []
        
        for file_info in files:
            filepath = file_info.get("path", "")
            content = file_info.get("content", "")
            
            if not filepath:
                continue
            
            full_path = project_path / filepath
            full_path.parent.mkdir(parents=True, exist_ok=True)
            
            try:
                full_path.write_text(content, encoding="utf-8")
                relative_path = str(full_path.relative_to(workspace_path))
                created_files.append(relative_path)
            except Exception as e:
                errors.append(f"{filepath}: {e}")
        
        if created_files:
            return ToolResult(
                success=True,
                message=f"项目已创建: {project_name}，共 {len(created_files)} 个文件",
                files_created=created_files,
                data={
                    "project_name": project_name,
                    "total_files": len(created_files),
                    "errors": errors if errors else None,
                }
            )
        else:
            return ToolResult(
                success=False,
                message="没有文件被创建",
                error="; ".join(errors) if errors else "files list is empty",
            )


class GeneratePPTTool(BaseTool):
    """
    生成 PPT 工具
    """
    
    name = "generate_ppt"
    display_name = "生成PPT"
    description = """根据内容生成 PowerPoint 演示文稿。
    
使用场景：
- 产品发布会 PPT
- 工作汇报 PPT
- 培训演示 PPT
- 商业计划书 PPT

参数说明：
- title: PPT 标题
- slides: 幻灯片内容列表，每页包含 title 和 points
"""
    category = "document"
    
    parameters = [
        ToolParameter(
            name="title",
            description="PPT 标题",
            type="string",
            required=True,
        ),
        ToolParameter(
            name="slides",
            description="幻灯片列表，格式: [{\"title\": \"页面标题\", \"points\": [\"要点1\", \"要点2\"]}]",
            type="array",
            required=True,
        ),
        ToolParameter(
            name="filename",
            description="文件名（不含扩展名）",
            type="string",
            required=False,
        ),
    ]
    
    async def execute(
        self,
        title: str,
        slides: List[dict],
        filename: Optional[str] = None,
    ) -> ToolResult:
        """生成 PPT"""
        if not self.workspace_id:
            return ToolResult(
                success=False,
                message="未指定工作空间",
                error="workspace_id is required",
            )
        
        workspace_path = WORKSPACES_ROOT / self.workspace_id
        if not workspace_path.exists():
            return ToolResult(
                success=False,
                message=f"工作空间不存在: {self.workspace_id}",
                error="workspace not found",
            )
        
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except ImportError:
            return ToolResult(
                success=False,
                message="缺少 python-pptx 库",
                error="请安装: pip install python-pptx",
                suggestions=["运行: pip install python-pptx"]
            )
        
        # 创建 PPT
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        # 封面
        title_slide_layout = prs.slide_layouts[6]  # 空白布局
        slide = prs.slides.add_slide(title_slide_layout)
        
        # 添加标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.33), Inches(1.5))
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.alignment = 1  # 居中
        
        # 内容页
        for slide_data in slides:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            # 页面标题
            slide_title = slide_data.get("title", "")
            if slide_title:
                title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8))
                tf = title_box.text_frame
                tf.paragraphs[0].text = slide_title
                tf.paragraphs[0].font.size = Pt(32)
                tf.paragraphs[0].font.bold = True
            
            # 要点
            points = slide_data.get("points", [])
            if points:
                content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(11.93), Inches(5.5))
                tf = content_box.text_frame
                tf.word_wrap = True
                
                for i, point in enumerate(points):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = f"• {point}"
                    p.font.size = Pt(20)
                    p.space_before = Pt(12)
        
        # 保存
        if not filename:
            timestamp = datetime.now().strftime("%H%M%S")
            filename = f"{title[:20]}_{timestamp}"
        
        ppt_dir = workspace_path / "ppt"
        ppt_dir.mkdir(parents=True, exist_ok=True)
        
        filepath = ppt_dir / f"{filename}.pptx"
        prs.save(str(filepath))
        
        relative_path = str(filepath.relative_to(workspace_path))
        
        return ToolResult(
            success=True,
            message=f"PPT 已生成: {relative_path}",
            files_created=[relative_path],
            data={
                "filepath": relative_path,
                "slide_count": len(slides) + 1,
            }
        )


class GenerateExcelTool(BaseTool):
    """
    生成 Excel 工具
    """
    
    name = "generate_excel"
    display_name = "生成Excel"
    description = """生成 Excel 表格文件。
    
使用场景：
- 数据表格
- 财务报表
- 统计分析
- 工资表

参数说明：
- title: 表格标题
- headers: 表头列表
- rows: 数据行列表
"""
    category = "document"
    
    parameters = [
        ToolParameter(
            name="title",
            description="表格标题",
            type="string",
            required=True,
        ),
        ToolParameter(
            name="headers",
            description="表头列表，如 [\"姓名\", \"部门\", \"工资\"]",
            type="array",
            required=True,
        ),
        ToolParameter(
            name="rows",
            description="数据行列表，如 [[\"张三\", \"技术部\", 10000], [\"李四\", \"销售部\", 8000]]",
            type="array",
            required=True,
        ),
        ToolParameter(
            name="filename",
            description="文件名（不含扩展名）",
            type="string",
            required=False,
        ),
    ]
    
    async def execute(
        self,
        title: str,
        headers: List[str],
        rows: List[List],
        filename: Optional[str] = None,
    ) -> ToolResult:
        """生成 Excel"""
        if not self.workspace_id:
            return ToolResult(
                success=False,
                message="未指定工作空间",
                error="workspace_id is required",
            )
        
        workspace_path = WORKSPACES_ROOT / self.workspace_id
        if not workspace_path.exists():
            return ToolResult(
                success=False,
                message=f"工作空间不存在: {self.workspace_id}",
                error="workspace not found",
            )
        
        try:
            from openpyxl import Workbook
            from openpyxl.styles import Font, Alignment, PatternFill
        except ImportError:
            return ToolResult(
                success=False,
                message="缺少 openpyxl 库",
                error="请安装: pip install openpyxl",
                suggestions=["运行: pip install openpyxl"]
            )
        
        # 创建工作簿
        wb = Workbook()
        ws = wb.active
        ws.title = title[:31]  # Excel sheet name 最长 31 字符
        
        # 标题行样式
        header_font = Font(bold=True, color="FFFFFF")
        header_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
        
        # 写入表头
        for col, header in enumerate(headers, 1):
            cell = ws.cell(row=1, column=col, value=header)
            cell.font = header_font
            cell.fill = header_fill
            cell.alignment = Alignment(horizontal="center")
        
        # 写入数据
        for row_idx, row_data in enumerate(rows, 2):
            for col_idx, value in enumerate(row_data, 1):
                ws.cell(row=row_idx, column=col_idx, value=value)
        
        # 自动调整列宽
        for col in ws.columns:
            max_length = 0
            column = col[0].column_letter
            for cell in col:
                try:
                    if len(str(cell.value)) > max_length:
                        max_length = len(str(cell.value))
                except:
                    pass
            adjusted_width = min(max_length + 2, 50)
            ws.column_dimensions[column].width = adjusted_width
        
        # 保存
        if not filename:
            timestamp = datetime.now().strftime("%H%M%S")
            filename = f"{title[:20]}_{timestamp}"
        
        data_dir = workspace_path / "data"
        data_dir.mkdir(parents=True, exist_ok=True)
        
        filepath = data_dir / f"{filename}.xlsx"
        wb.save(str(filepath))
        
        relative_path = str(filepath.relative_to(workspace_path))
        
        return ToolResult(
            success=True,
            message=f"Excel 已生成: {relative_path}",
            files_created=[relative_path],
            data={
                "filepath": relative_path,
                "row_count": len(rows),
                "column_count": len(headers),
            }
        )
